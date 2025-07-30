import os
import json
import uuid
import zipfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from agents import PPTCrew
from config import Config
from themes import PPTThemes, ThemeConfig
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTProjectManager:
    """
    Manages the complete PPT generation process from user input to final deliverable.
    Coordinates AI agents and handles file generation.
    """
    
    def __init__(self, socketio=None):
        self.socketio = socketio
        self.crew = PPTCrew()
        self.projects = {}  # Store project information
        self.state_file = os.path.join(Config.TEMP_DIR, 'project_states.json')
        
        # Ensure directories exist
        os.makedirs(Config.GENERATED_PPTS_DIR, exist_ok=True)
        os.makedirs(Config.TEMP_DIR, exist_ok=True)
        
        # Load persisted project states
        self._load_project_states()
    
    def _save_project_state(self, project_id):
        """Save project state to disk to survive restarts."""
        try:
            # Only save essential state information
            state_data = {}
            if project_id in self.projects:
                project = self.projects[project_id].copy()
                # Remove non-serializable objects
                project.pop('socketio', None)
                state_data[project_id] = project
            
            # Load existing states and update
            existing_states = {}
            if os.path.exists(self.state_file):
                try:
                    with open(self.state_file, 'r') as f:
                        existing_states = json.load(f)
                except (json.JSONDecodeError, IOError):
                    pass
            
            existing_states.update(state_data)
            
            with open(self.state_file, 'w') as f:
                json.dump(existing_states, f, indent=2)
                
        except Exception as e:
            logger.warning(f"Failed to save project state: {e}")
    
    def _validate_plan_data(self, plan_data):
        """
        Validate and fix plan_data structure to ensure it has all required fields.
        """
        if not isinstance(plan_data, dict):
            logger.error(f"plan_data is not a dictionary: {type(plan_data)}")
            return False
        
        # Ensure required top-level keys exist
        required_keys = ['presentation_title', 'presentation_description', 'slides']
        for key in required_keys:
            if key not in plan_data:
                logger.warning(f"Missing required key '{key}' in plan_data")
                if key == 'presentation_title':
                    plan_data[key] = 'Generated Presentation'
                elif key == 'presentation_description':
                    plan_data[key] = 'Created with AI'
                elif key == 'slides':
                    plan_data[key] = []
        
        # Ensure slides is a list
        if not isinstance(plan_data.get('slides'), list):
            logger.warning("'slides' is not a list, converting or creating empty list")
            plan_data['slides'] = []
        
        # Validate each slide
        for i, slide in enumerate(plan_data['slides']):
            if not isinstance(slide, dict):
                logger.warning(f"Slide {i} is not a dictionary: {type(slide)}")
                plan_data['slides'][i] = {
                    'title': f'Slide {i + 1}',
                    'content': 'Generated content',
                    'content_type': 'bullet_points',
                    'bullet_points': ['Point 1', 'Point 2', 'Point 3']
                }
                continue
            
            # Ensure required slide keys
            slide_defaults = {
                'title': f'Slide {i + 1}',
                'content': 'Generated content',
                'content_type': 'bullet_points',
                'bullet_points': ['Point 1', 'Point 2', 'Point 3'],
                'color_scheme': 'professional_blue'
            }
            
            for key, default_value in slide_defaults.items():
                if key not in slide:
                    slide[key] = default_value
            
            # Clean markdown formatting from slide content
            if 'title' in slide:
                slide['title'] = self._clean_markdown_formatting(slide['title'])
            
            if 'content' in slide:
                slide['content'] = self._clean_markdown_formatting(slide['content'])
            
            if 'bullet_points' in slide and isinstance(slide['bullet_points'], list):
                slide['bullet_points'] = [
                    self._clean_markdown_formatting(str(point)) 
                    for point in slide['bullet_points']
                ]
            
            if 'numbered_points' in slide and isinstance(slide['numbered_points'], list):
                slide['numbered_points'] = [
                    self._clean_markdown_formatting(str(point)) 
                    for point in slide['numbered_points']
                ]
        
        # Clean markdown from title and description
        plan_data['presentation_title'] = self._clean_markdown_formatting(plan_data['presentation_title'])
        plan_data['presentation_description'] = self._clean_markdown_formatting(plan_data['presentation_description'])
        
        return True
    
    def _clean_json_content(self, content):
        """
        Clean JSON content that might be wrapped in markdown code blocks.
        CrewAI often returns JSON wrapped in ```json ... ``` blocks.
        """
        if not isinstance(content, str):
            return content
        
        # Remove markdown code block formatting
        content = content.strip()
        
        # Remove ```json at the beginning
        if content.startswith('```json'):
            content = content[7:]  # Remove '```json'
        elif content.startswith('```'):
            content = content[3:]   # Remove '```'
        
        # Remove ``` at the end
        if content.endswith('```'):
            content = content[:-3]
        
        # Strip any remaining whitespace
        content = content.strip()
        
        return content
    
    def _clean_markdown_formatting(self, text):
        """
        Clean markdown formatting from text content.
        Removes **bold**, *italic*, and other markdown syntax.
        """
        if not isinstance(text, str):
            return text
        
        import re
        
        # Remove bold formatting (**text** or __text__)
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'__(.*?)__', r'\1', text)
        
        # Remove italic formatting (*text* or _text_)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        text = re.sub(r'_(.*?)_', r'\1', text)
        
        # Remove code formatting (`text`)
        text = re.sub(r'`(.*?)`', r'\1', text)
        
        # Remove strikethrough (~~text~~)
        text = re.sub(r'~~(.*?)~~', r'\1', text)
        
        # Remove bullet points and numbered lists at the beginning of lines
        text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
        
        # Clean up extra whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()
        
        return text
    
    def _extract_crew_result(self, crew_output):
        """
        Extract the actual result from a CrewOutput object.
        CrewAI returns a special CrewOutput object that needs to be handled properly.
        """
        logger.info(f"CrewOutput type: {type(crew_output)}")
        
        # Try to get available attributes (filter out private ones for cleaner output)
        available_attrs = [attr for attr in dir(crew_output) if not attr.startswith('_')]
        logger.info(f"CrewOutput public attributes: {available_attrs}")
        
        try:
            # Try different ways to extract the result in order of preference
            extraction_methods = [
                ('raw', lambda x: x.raw),
                ('result', lambda x: x.result),
                ('output', lambda x: x.output),
                ('text', lambda x: x.text),
                ('content', lambda x: x.content),
                ('str', lambda x: str(x))
            ]
            
            for method_name, method_func in extraction_methods:
                try:
                    if method_name == 'str' or hasattr(crew_output, method_name.replace('str', '__str__')):
                        result = method_func(crew_output)
                        logger.info(f"Successfully extracted using method: {method_name}")
                        logger.info(f"Extracted result type: {type(result)}")
                        if isinstance(result, str) and len(result) > 200:
                            logger.info(f"Result preview: {result[:200]}...")
                        else:
                            logger.info(f"Result: {result}")
                        return result
                except AttributeError:
                    continue
                except Exception as e:
                    logger.warning(f"Method {method_name} failed: {e}")
                    continue
            
            # If all methods fail, return a string representation
            logger.warning("All extraction methods failed, returning string representation")
            return str(crew_output)
                
        except Exception as e:
            logger.error(f"Failed to extract result from CrewOutput: {e}")
            return f"Error extracting result: {str(e)}"
    
    def _load_project_states(self):
        """Load project states from disk."""
        try:
            if os.path.exists(self.state_file):
                with open(self.state_file, 'r') as f:
                    self.projects = json.load(f)
                logger.info(f"Loaded {len(self.projects)} project states from disk")
        except Exception as e:
            logger.warning(f"Failed to load project states: {e}")
            self.projects = {}
        os.makedirs(Config.TEMP_DIR, exist_ok=True)
    
    def emit_progress(self, project_id, stage, message):
        """Emit progress updates via WebSocket and persist state"""
        # Update project stage in memory
        if project_id in self.projects:
            self.projects[project_id]['current_stage'] = stage
            self.projects[project_id]['last_message'] = message
            self.projects[project_id]['last_update'] = datetime.now().isoformat()
        
        if self.socketio:
            self.socketio.emit('progress_update', {
                'project_id': project_id,
                'stage': stage,
                'message': message,
                'timestamp': datetime.now().isoformat()
            }, room=project_id)
        
        # Persist state after each progress update
        self._save_project_state(project_id)
        logger.info(f"[{project_id}] {stage}: {message}")
    
    def generate_presentation(self, user_prompt, num_slides=5, project_id=None, theme_name='corporate_blue'):
        """
        Main method to generate a complete presentation.
        
        Args:
            user_prompt: User's description of the presentation
            num_slides: Number of slides to generate (default: 5)
            project_id: Optional project ID (will generate if not provided)
            theme_name: Theme to use for the presentation (default: 'corporate_blue')
        """
        if not project_id:
            project_id = str(uuid.uuid4())
        
        try:
            # Initialize project
            self.projects[project_id] = {
                'id': project_id,
                'prompt': user_prompt,
                'num_slides': num_slides,
                'theme_name': theme_name,
                'status': 'started',
                'created_at': datetime.now().isoformat(),
                'stages': []
            }
            
            self.emit_progress(project_id, 'initialization', 'Starting presentation generation...')
            
            # Stage 1: Planning
            self.emit_progress(project_id, 'planning', 'Planning Agent is analyzing your requirements...')
            
            # Stage 2: Content Creation
            self.emit_progress(project_id, 'content_creation', 'Content Creator is generating slide content...')
            
            # Stage 3: Design
            self.emit_progress(project_id, 'design', 'Designer is creating the visual presentation...')
            
            # Execute the complete CrewAI pipeline
            presentation_plan = self.crew.create_presentation_plan(user_prompt, num_slides)
            
            # Get the selected theme
            theme = PPTThemes.get_theme(theme_name)
            if not theme:
                logger.warning(f"Theme '{theme_name}' not found, using default corporate_blue")
                theme = PPTThemes.get_theme('corporate_blue')
            
            logger.info(f"Using theme: {theme.display_name}")
            
            # Parse the result - CrewOutput object needs special handling
            try:
                # Extract the actual result from CrewOutput object
                raw_result = self._extract_crew_result(presentation_plan)
                
                # Log the raw result for debugging
                logger.info(f"Raw result type: {type(raw_result)}")
                if isinstance(raw_result, str):
                    logger.info(f"Raw result preview: {raw_result[:1000]}...")
                
                # Now try to parse the extracted result
                if isinstance(raw_result, str):
                    # Clean any markdown formatting from the JSON
                    cleaned_result = self._clean_json_content(raw_result)
                    
                    # Log cleaned result for debugging
                    logger.info(f"Cleaned result preview: {cleaned_result[:1000]}...")
                    
                    # Try to parse as JSON
                    try:
                        plan_data = json.loads(cleaned_result)
                        logger.info("Successfully parsed CrewOutput as JSON after cleaning")
                        logger.info(f"Plan data keys: {list(plan_data.keys()) if isinstance(plan_data, dict) else 'Not a dict'}")
                        
                        # Log sample slide content to verify content generation
                        if isinstance(plan_data, dict) and 'slides' in plan_data and plan_data['slides']:
                            sample_slide = plan_data['slides'][0]
                            logger.info(f"Sample slide keys: {list(sample_slide.keys()) if isinstance(sample_slide, dict) else 'Not a dict'}")
                            if isinstance(sample_slide, dict):
                                logger.info(f"Sample slide content_type: {sample_slide.get('content_type', 'missing')}")
                                logger.info(f"Sample slide has content: {'content' in sample_slide}")
                                logger.info(f"Sample slide has bullet_points: {'bullet_points' in sample_slide}")
                        
                    except json.JSONDecodeError as json_err:
                        # If not valid JSON, create fallback plan
                        logger.warning(f"CrewOutput result is not valid JSON after cleaning: {json_err}")
                        logger.debug(f"Cleaned result was: {cleaned_result[:500]}...")
                        plan_data = self._create_fallback_plan(user_prompt, num_slides)
                elif isinstance(raw_result, dict):
                    # Already a dictionary
                    plan_data = raw_result
                    logger.info("CrewOutput was already a dictionary")
                    
                    # Log sample slide content to verify content generation
                    if 'slides' in plan_data and plan_data['slides']:
                        sample_slide = plan_data['slides'][0]
                        logger.info(f"Sample slide keys: {list(sample_slide.keys()) if isinstance(sample_slide, dict) else 'Not a dict'}")
                        if isinstance(sample_slide, dict):
                            logger.info(f"Sample slide content_type: {sample_slide.get('content_type', 'missing')}")
                            logger.info(f"Sample slide has content: {'content' in sample_slide}")
                            logger.info(f"Sample slide has bullet_points: {'bullet_points' in sample_slide}")
                            
                else:
                    # Unknown format, use fallback
                    logger.warning(f"Unknown CrewOutput format: {type(raw_result)}, using fallback plan")
                    plan_data = self._create_fallback_plan(user_prompt, num_slides)
                    
            except Exception as e:
                # If all parsing fails, create a basic structure
                logger.error(f"Failed to parse CrewOutput: {e}, using fallback plan")
                plan_data = self._create_fallback_plan(user_prompt, num_slides)
            
            # Validate and fix plan_data structure
            logger.info(f"Validating plan data structure... Type: {type(plan_data)}")
            
            # Additional check: if plan_data is still a CrewOutput, extract it
            if hasattr(plan_data, 'raw') or str(type(plan_data)).lower().find('crewoutput') != -1:
                logger.warning("plan_data is still a CrewOutput object, extracting content...")
                plan_data = self._extract_crew_result(plan_data)
                
                if isinstance(plan_data, str):
                    cleaned_result = self._clean_json_content(plan_data)
                    try:
                        plan_data = json.loads(cleaned_result)
                        logger.info("Successfully parsed nested CrewOutput as JSON")
                    except json.JSONDecodeError:
                        logger.warning("Nested CrewOutput is not valid JSON, using fallback")
                        plan_data = self._create_fallback_plan(user_prompt, num_slides)
            
            if not self._validate_plan_data(plan_data):
                logger.warning("Plan data validation failed, using fallback plan")
                plan_data = self._create_fallback_plan(user_prompt, num_slides)
                self._validate_plan_data(plan_data)  # Ensure fallback is also valid
            
            logger.info(f"Final plan_data type before PowerPoint creation: {type(plan_data)}")
            logger.info(f"Plan_data keys: {list(plan_data.keys()) if isinstance(plan_data, dict) else 'Not a dictionary'}")
            
            # Stage 4: Generate PowerPoint
            self.emit_progress(project_id, 'packaging', 'Creating the PowerPoint presentation...')
            ppt_path = self._create_powerpoint(plan_data, project_id, theme)
            
            # Stage 5: Finalization
            self.emit_progress(project_id, 'finalization', 'Finalizing your presentation...')
            
            # Update project with results
            self.projects[project_id].update({
                'status': 'completed',
                'plan': plan_data,
                'ppt_path': ppt_path,
                'completed_at': datetime.now().isoformat()
            })
            
            self.emit_progress(project_id, 'completed', 'Presentation generated successfully!')
            
            if self.socketio:
                self.socketio.emit('project_completed', {
                    'project_id': project_id,
                    'result': {
                        'success': True,
                        'presentation_plan': plan_data,
                        'file_path': ppt_path
                    }
                }, room=project_id)
            
            return {
                'success': True,
                'project_id': project_id,
                'plan': plan_data,
                'file_path': ppt_path
            }
            
        except Exception as e:
            error_str = str(e)
            logger.error(f"Error generating presentation: {error_str}")
            logger.error(f"Error type: {type(e)}")
            
            # Provide specific error messages for common issues
            if "CrewOutput" in error_str and "get" in error_str:
                user_message = "There was an issue processing the AI-generated content. The system will use a fallback approach."
                self.emit_progress(project_id, 'failed', user_message)
            elif "503" in error_str or "overloaded" in error_str.lower() or "unavailable" in error_str.lower():
                user_message = "The AI service is currently overloaded. This is a temporary issue. Please try again in a few minutes."
                self.emit_progress(project_id, 'failed', user_message)
            elif "api key" in error_str.lower() and ("invalid" in error_str.lower() or "expired" in error_str.lower()):
                user_message = "API key issue detected. Please check your API key configuration."
                self.emit_progress(project_id, 'failed', user_message)
            else:
                user_message = f"An error occurred during generation: {error_str}"
                self.emit_progress(project_id, 'failed', user_message)
            
            self.projects[project_id]['status'] = 'failed'
            self.projects[project_id]['error'] = error_str
            self.projects[project_id]['user_message'] = user_message
            
            if self.socketio:
                self.socketio.emit('project_failed', {
                    'project_id': project_id,
                    'error': error_str,
                    'user_message': user_message
                }, room=project_id)
            
            return {
                'success': False,
                'project_id': project_id,
                'error': error_str,
                'user_message': user_message
            }
    
    def _create_fallback_plan(self, user_prompt, num_slides):
        """Create a basic presentation plan with varied slide types if AI generation fails"""
        content_types = ['paragraph', 'bullet_points', 'numbered_list', 'two_column']
        
        slides = []
        for i in range(num_slides):
            content_type = content_types[i % len(content_types)]
            
            slide = {
                "slide_number": i + 1,
                "title": f"Key Topic {i + 1}: {user_prompt[:30]}",
                "description": f"Content for slide {i + 1}",
                "content_type": content_type,
                "layout_style": "standard",
                "color_scheme": "professional_blue"
            }
            
            if content_type == "bullet_points":
                slide["bullet_points"] = [
                    f"Important insight about {user_prompt[:25]} that provides value",
                    f"Key benefit or feature related to {user_prompt[:20]} implementation",
                    f"Strategic advantage gained through {user_prompt[:20]} adoption",
                    f"Measurable outcome expected from {user_prompt[:20]} process"
                ]
            elif content_type == "paragraph":
                slide["content"] = f"This section explores the fundamental aspects of {user_prompt}. The implementation requires careful consideration of multiple factors including strategic alignment, resource allocation, and stakeholder engagement. By focusing on these core elements, organizations can achieve sustainable growth and competitive advantage in their respective markets."
            elif content_type == "numbered_list":
                slide["numbered_points"] = [
                    f"Initial assessment and planning phase for {user_prompt[:20]}",
                    f"Implementation strategy development and resource allocation",
                    f"Execution phase with milestone tracking and quality control",
                    f"Review and optimization based on performance metrics"
                ]
            elif content_type == "two_column":
                slide["left_content"] = f"Current State Analysis: The existing situation regarding {user_prompt} shows various challenges and opportunities that need to be addressed systematically."
                slide["right_content"] = f"Future Vision: The desired outcome for {user_prompt} includes improved efficiency, enhanced performance, and measurable business value."
            
            slides.append(slide)
        
        return {
            "presentation_title": f"Comprehensive Overview: {user_prompt[:40]}",
            "presentation_description": f"A detailed analysis and strategic guide for {user_prompt}",
            "total_slides": num_slides,
            "slides": slides
        }
    
    def _create_powerpoint(self, plan_data, project_id, theme: ThemeConfig):
        """
        Create the actual PowerPoint file from the plan data using the specified theme.
        """
        logger.info(f"Creating PowerPoint with theme: {theme.display_name}")
        logger.info(f"Plan_data type: {type(plan_data)}")
        
        # Safety check: if plan_data is still a CrewOutput, extract and parse it
        if hasattr(plan_data, 'raw') or str(type(plan_data)).lower().find('crewoutput') != -1:
            logger.warning("plan_data is a CrewOutput in _create_powerpoint, extracting...")
            raw_result = self._extract_crew_result(plan_data)
            
            if isinstance(raw_result, str):
                cleaned_result = self._clean_json_content(raw_result)
                try:
                    plan_data = json.loads(cleaned_result)
                    logger.info("Successfully extracted and parsed CrewOutput in PowerPoint creation")
                except json.JSONDecodeError:
                    logger.error("Failed to parse CrewOutput in PowerPoint creation, using fallback")
                    plan_data = self._create_fallback_plan("Generated Presentation", 5)
            elif isinstance(raw_result, dict):
                plan_data = raw_result
            else:
                logger.error("Unknown CrewOutput format in PowerPoint creation, using fallback")
                plan_data = self._create_fallback_plan("Generated Presentation", 5)
        
        # Ensure plan_data is a dictionary
        if not isinstance(plan_data, dict):
            logger.error(f"plan_data is not a dictionary: {type(plan_data)}, using fallback")
            plan_data = self._create_fallback_plan("Generated Presentation", 5)
        
        logger.info(f"Final plan_data type in PowerPoint creation: {type(plan_data)}")
        logger.info(f"Available keys: {list(plan_data.keys()) if isinstance(plan_data, dict) else 'None'}")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Apply professional background to title slide
        self._apply_professional_background(slide, theme, is_title_slide=True, slide_type="title")
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self._clean_markdown_formatting(plan_data.get('presentation_title', 'Generated Presentation'))
        subtitle.text = self._clean_markdown_formatting(plan_data.get('presentation_description', 'Created with AI'))
        
        # Style title slide with theme
        self._style_title_slide(title, subtitle, theme)
        
        # Content slides
        for slide_data in plan_data.get('slides', []):
            self._create_content_slide(prs, slide_data, theme)
        
        # Save presentation
        filename = f"presentation_{project_id}.pptx"
        file_path = os.path.join(Config.GENERATED_PPTS_DIR, filename)
        prs.save(file_path)
        
        return file_path
    
    def _style_title_slide(self, title, subtitle, theme: ThemeConfig):
        """Apply professional styling to title slide with better alignment using theme colors"""
        # Title styling - work with existing text frame
        try:
            title_frame = title.text_frame
            title_p = title_frame.paragraphs[0]
            
            # Set title text if not already set
            if not title_p.text:
                title_p.text = title.text if hasattr(title, 'text') else "Presentation Title"
            
            title_font = title_p.font
            title_font.name = theme.font_scheme.title_font
            title_font.size = Pt(theme.font_scheme.title_size + 18)  # Larger for title slide
            title_font.bold = True
            title_font.color.rgb = theme.color_scheme.text_light  # White text for title slides
            
            # Center align title
            title_p.alignment = PP_ALIGN.CENTER
            
            # Set margins safely
            try:
                title_frame.margin_top = Inches(1.5)
                title_frame.margin_bottom = Inches(0.5)
                title_frame.margin_left = Inches(0.5)
                title_frame.margin_right = Inches(0.5)
            except:
                pass  # If margin setting fails, continue
            
        except Exception as e:
            logger.warning(f"Error styling title: {e}")
        
        # Subtitle styling - work with existing text frame
        try:
            subtitle_frame = subtitle.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            
            # Set subtitle text if not already set
            if not subtitle_p.text:
                subtitle_p.text = subtitle.text if hasattr(subtitle, 'text') else "Subtitle"
            
            subtitle_font = subtitle_p.font
            subtitle_font.name = theme.font_scheme.content_font
            subtitle_font.size = Pt(theme.font_scheme.subtitle_size)
            subtitle_font.color.rgb = theme.color_scheme.text_secondary if theme.name != 'elegant_dark' and theme.name != 'tech_cyber' else RGBColor(230, 230, 230)
            
            # Center align subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            
            # Set margins safely
            try:
                subtitle_frame.margin_top = Inches(0.3)
                subtitle_frame.margin_bottom = Inches(1.0)
                subtitle_frame.margin_left = Inches(1.0)
                subtitle_frame.margin_right = Inches(1.0)
            except:
                pass  # If margin setting fails, continue
                
        except Exception as e:
            logger.warning(f"Error styling subtitle: {e}")
        
        # Position title and subtitle better (optional enhancement)
        try:
            title.top = Inches(2.5)
            title.left = Inches(0.5)
            title.width = Inches(12.33)
            title.height = Inches(1.5)
            
            subtitle.top = Inches(4.2)
            subtitle.left = Inches(1.0)
            subtitle.width = Inches(11.33)
            subtitle.height = Inches(1.0)
        except Exception as e:
            logger.debug(f"Could not adjust title/subtitle positioning: {e}")
            # Continue with default positioning
    
    def _apply_professional_background(self, slide, theme: ThemeConfig, is_title_slide=False, slide_type="content"):
        """Apply creative professional gradient background to slide using theme colors"""
        background = slide.background
        fill = background.fill
        
        try:
            # Apply gradient background
            fill.gradient()
            fill.gradient_angle = theme.gradient_angle
            
            # Get gradient stops and set colors
            gradient_stops = fill.gradient_stops
            
            # Set the first stop (position 0.0)
            if len(gradient_stops) > 0:
                if is_title_slide:
                    # Title slides use primary colors for more impact
                    gradient_stops[0].color.rgb = theme.color_scheme.primary
                else:
                    # Content slides use background colors
                    gradient_stops[0].color.rgb = theme.color_scheme.background_start
                gradient_stops[0].position = 0.0
            
            # Set the second stop (position 1.0) - create if doesn't exist
            if len(gradient_stops) > 1:
                if is_title_slide:
                    gradient_stops[1].color.rgb = theme.color_scheme.secondary
                else:
                    gradient_stops[1].color.rgb = theme.color_scheme.background_end
                gradient_stops[1].position = 1.0
            else:
                # Add second stop if it doesn't exist
                try:
                    stop = gradient_stops.add_gradient_stop(1.0)
                    if is_title_slide:
                        stop.color.rgb = theme.color_scheme.secondary
                    else:
                        stop.color.rgb = theme.color_scheme.background_end
                except:
                    # If add_gradient_stop doesn't work, use simple solid fill
                    fill.solid()
                    if is_title_slide:
                        fill.fore_color.rgb = theme.color_scheme.primary
                    else:
                        fill.fore_color.rgb = theme.color_scheme.background_start
                        
        except Exception as e:
            logger.warning(f"Error applying gradient background: {e}")
            # Fallback to solid color
            try:
                fill.solid()
                if is_title_slide:
                    fill.fore_color.rgb = theme.color_scheme.primary
                else:
                    fill.fore_color.rgb = theme.color_scheme.background_start
            except Exception as fallback_e:
                logger.error(f"Error applying fallback background: {fallback_e}")
                # Continue with default background
    
    def _create_content_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a content slide with varied layouts based on content type using theme styling"""
        content_type = slide_data.get('content_type', 'bullet_points')
        layout_style = slide_data.get('layout_style', 'standard')
        
        # Choose different layouts based on content
        if content_type == 'title_only':
            self._create_title_only_slide(prs, slide_data, theme)
        elif content_type == 'two_column':
            self._create_two_column_slide(prs, slide_data, theme)
        elif content_type == 'image_focus':
            self._create_image_focus_slide(prs, slide_data, theme)
        elif content_type == 'comparison':
            self._create_comparison_slide(prs, slide_data, theme)
        else:
            self._create_standard_content_slide(prs, slide_data, theme)
    
    def _create_standard_content_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a standard content slide with enhanced formatting using theme styling"""
        # Use content slide layout
        content_slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Apply creative background with theme
        self._apply_professional_background(slide, theme, is_title_slide=False, slide_type="content")
        
        # Set title
        title = slide.shapes.title
        title.text = self._clean_markdown_formatting(slide_data.get('title', 'Slide Title'))
        
        # Style title with theme
        title_font = title.text_frame.paragraphs[0].font
        title_font.name = theme.font_scheme.title_font
        title_font.size = Pt(theme.font_scheme.title_size)
        title_font.bold = True
        title_font.color.rgb = theme.color_scheme.primary
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add content
        content_placeholder = slide.placeholders[1]
        content_frame = content_placeholder.text_frame
        content_frame.clear()
        
        # Set content frame properties
        content_frame.word_wrap = True
        content_frame.auto_size = None
        
        content_type = slide_data.get('content_type', 'bullet_points')
        
        if content_type == 'bullet_points' and slide_data.get('bullet_points'):
            self._add_bullet_points_with_icons(content_frame, slide_data['bullet_points'], theme)
        
        elif content_type == 'paragraph' and slide_data.get('content'):
            self._add_paragraph_content(content_frame, slide_data['content'], theme)
        
        elif content_type == 'numbered_list' and slide_data.get('numbered_points'):
            self._add_numbered_list(content_frame, slide_data['numbered_points'], theme)
        
        else:
            # Fallback content
            p = content_frame.paragraphs[0]
            p.text = "Content will be added here"
            p.font.name = theme.font_scheme.content_font
            p.font.size = Pt(theme.font_scheme.content_size)
            p.font.color.rgb = theme.color_scheme.text_secondary
            p.alignment = PP_ALIGN.CENTER
    
    def _add_bullet_points_with_icons(self, content_frame, bullet_points, theme: ThemeConfig):
        """Add bullet points with relevant icons where appropriate using theme styling"""
        # Icon mappings for common concepts - enhanced based on theme style
        if theme.icon_style == 'modern':
            icon_mappings = {
                'success': '✅', 'achieve': '🎯', 'goal': '🎯', 'target': '🎯',
                'complete': '✅', 'finish': '✅', 'done': '✅',
                'growth': '📈', 'increase': '📈', 'improve': '�', 'rise': '📈',
                'expand': '📈', 'scale': '📈', 'develop': '�',
                'idea': '💡', 'innovation': '⚡', 'creative': '💡', 'solution': '�',
                'insight': '💡', 'concept': '💡',
                'time': '⏰', 'schedule': '📅', 'deadline': '⏰', 'timeline': '📅',
                'money': '💰', 'cost': '💰', 'budget': '💰', 'revenue': '💰',
                'communicate': '💬', 'message': '�', 'feedback': '💬', 'discuss': '💬',
                'technology': '⚡', 'digital': '💻', 'software': '�', 'system': '⚙️',
                'security': '�', 'protect': '�️', 'safe': '🔒', 'risk': '⚠️',
                'quality': '⭐', 'excellent': '⭐', 'best': '🏆', 'premium': '⭐'
            }
        elif theme.icon_style == 'creative':
            icon_mappings = {
                'success': '🌟', 'achieve': '🎯', 'goal': '🎯', 'target': '🎯',
                'complete': '🌟', 'finish': '🌟', 'done': '🌟',
                'growth': '🌱', 'increase': '�', 'improve': '🌱', 'rise': '�',
                'expand': '🌱', 'scale': '🌱', 'develop': '🌱',
                'idea': '�', 'innovation': '🎨', 'creative': '🎨', 'solution': '�',
                'insight': '💡', 'concept': '�',
                'time': '⏰', 'schedule': '📅', 'deadline': '⏰', 'timeline': '📅',
                'money': '💎', 'cost': '💎', 'budget': '💎', 'revenue': '💎',
                'communicate': '💬', 'message': '💬', 'feedback': '💬', 'discuss': '💬',
                'technology': '🔮', 'digital': '💻', 'software': '�', 'system': '⚙️',
                'security': '�️', 'protect': '🛡️', 'safe': '�️', 'risk': '⚠️',
                'quality': '💎', 'excellent': '💎', 'best': '🏆', 'premium': '💎'
            }
        elif theme.icon_style == 'minimal':
            icon_mappings = {
                'success': '✓', 'achieve': '→', 'goal': '→', 'target': '→',
                'complete': '✓', 'finish': '✓', 'done': '✓',
                'growth': '↗', 'increase': '↗', 'improve': '↗', 'rise': '↗',
                'expand': '↗', 'scale': '↗', 'develop': '↗',
                'idea': '•', 'innovation': '•', 'creative': '•', 'solution': '•',
                'insight': '•', 'concept': '•',
                'time': '•', 'schedule': '•', 'deadline': '•', 'timeline': '•',
                'money': '•', 'cost': '•', 'budget': '•', 'revenue': '•',
                'communicate': '•', 'message': '•', 'feedback': '•', 'discuss': '•',
                'technology': '•', 'digital': '•', 'software': '•', 'system': '•',
                'security': '•', 'protect': '•', 'safe': '•', 'risk': '!',
                'quality': '•', 'excellent': '•', 'best': '•', 'premium': '•'
            }
        else:  # classic
            icon_mappings = {
                'success': '✓', 'achieve': '•', 'goal': '•', 'target': '•',
                'complete': '✓', 'finish': '✓', 'done': '✓',
                'growth': '▲', 'increase': '▲', 'improve': '▲', 'rise': '▲',
                'expand': '▲', 'scale': '▲', 'develop': '▲',
                'idea': '◆', 'innovation': '◆', 'creative': '◆', 'solution': '◆',
                'insight': '◆', 'concept': '◆',
                'time': '●', 'schedule': '●', 'deadline': '●', 'timeline': '●',
                'money': '●', 'cost': '●', 'budget': '●', 'revenue': '●',
                'communicate': '●', 'message': '●', 'feedback': '●', 'discuss': '●',
                'technology': '●', 'digital': '●', 'software': '●', 'system': '●',
                'security': '●', 'protect': '●', 'safe': '●', 'risk': '▼',
                'quality': '●', 'excellent': '●', 'best': '●', 'premium': '●'
            }
        
        for i, point in enumerate(bullet_points):
            clean_point = self._clean_markdown_formatting(str(point))
            
            # Find appropriate icon
            icon = '•'  # Default bullet
            point_lower = clean_point.lower()
            for keyword, emoji in icon_mappings.items():
                if keyword in point_lower:
                    icon = emoji
                    break
            
            # Create paragraph
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"{icon} {clean_point}"
            p.level = 0
            
            # Style bullet points with theme
            p.font.name = theme.font_scheme.content_font
            p.font.size = Pt(theme.font_scheme.bullet_size)
            p.font.color.rgb = theme.color_scheme.text_primary
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)
    
    def _add_paragraph_content(self, content_frame, content, theme: ThemeConfig):
        """Add paragraph content with proper formatting using theme styling"""
        clean_content = self._clean_markdown_formatting(str(content))
        
        # Split into paragraphs if needed
        paragraphs = clean_content.split('\n\n')
        
        for i, paragraph_text in enumerate(paragraphs):
            if not paragraph_text.strip():
                continue
                
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = paragraph_text.strip()
            p.level = 0
            
            # Style paragraphs with theme
            p.font.name = theme.font_scheme.content_font
            p.font.size = Pt(theme.font_scheme.content_size)
            p.font.color.rgb = theme.color_scheme.text_primary
            p.alignment = PP_ALIGN.JUSTIFY  # Justify for better readability
            p.space_after = Pt(14)
    
    def _add_numbered_list(self, content_frame, numbered_points, theme: ThemeConfig):
        """Add numbered list with professional formatting using theme styling"""
        for i, point in enumerate(numbered_points):
            clean_point = self._clean_markdown_formatting(str(point))
            
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"{i + 1}. {clean_point}"
            p.level = 0
            
            # Style numbered points with theme
            p.font.name = theme.font_scheme.content_font
            p.font.size = Pt(theme.font_scheme.bullet_size)
            p.font.color.rgb = theme.color_scheme.text_primary
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)
    
    def _create_title_only_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a slide with just a large title (section divider)"""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply special background for title-only slides
        self._apply_professional_background(slide, theme, is_title_slide=True)
        
        # Add large centered title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2.5))
        title_frame = title_box.text_frame
        title_frame.clear()
        
        p = title_frame.paragraphs[0]
        p.text = self._clean_markdown_formatting(slide_data.get('title', 'Section Title'))
        p.alignment = PP_ALIGN.CENTER
        
        # Large, bold title font with theme
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = theme.color_scheme.text_light
    
    def _create_two_column_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a two-column layout slide"""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply background with theme
        self._apply_professional_background(slide, theme, is_title_slide=False)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = self._clean_markdown_formatting(slide_data.get('title', 'Two Column Layout'))
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(theme.font_scheme.title_size)
        p.font.bold = True
        p.font.color.rgb = theme.color_scheme.primary
        p.alignment = PP_ALIGN.CENTER
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
        left_frame = left_box.text_frame
        left_content = slide_data.get('left_content', slide_data.get('content', 'Left column content'))
        self._add_paragraph_content(left_frame, left_content, theme)
        
        # Right column  
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
        right_frame = right_box.text_frame
        right_content = slide_data.get('right_content', 'Right column content')
        if slide_data.get('bullet_points'):
            self._add_bullet_points_with_icons(right_frame, slide_data['bullet_points'][:3], theme)
        else:
            self._add_paragraph_content(right_frame, right_content, theme)
    
    def _create_comparison_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a comparison slide (Before vs After, Pros vs Cons, etc.)"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply background with theme
        self._apply_professional_background(slide, theme, is_title_slide=False)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = self._clean_markdown_formatting(slide_data.get('title', 'Comparison'))
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(theme.font_scheme.title_size)
        p.font.bold = True
        p.font.color.rgb = theme.color_scheme.primary
        p.alignment = PP_ALIGN.CENTER
        
        # Left side (Before/Pros) - use accent color for contrast
        left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5))
        left_title_frame = left_title.text_frame
        p = left_title_frame.paragraphs[0] 
        p.text = slide_data.get('left_title', 'Before')
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(24)
        p.font.color.rgb = theme.color_scheme.secondary
        p.alignment = PP_ALIGN.CENTER
        
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5))
        left_frame = left_box.text_frame
        if slide_data.get('left_points'):
            self._add_bullet_points_with_icons(left_frame, slide_data['left_points'], theme)
        
        # Right side (After/Cons) - use accent color for contrast
        right_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.5))
        right_title_frame = right_title.text_frame
        p = right_title_frame.paragraphs[0]
        p.text = slide_data.get('right_title', 'After')
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(24)
        p.font.color.rgb = theme.color_scheme.accent
        p.alignment = PP_ALIGN.CENTER
        
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5))
        right_frame = right_box.text_frame
        if slide_data.get('right_points'):
            self._add_bullet_points_with_icons(right_frame, slide_data['right_points'], theme)
    
    def _create_image_focus_slide(self, prs, slide_data, theme: ThemeConfig):
        """Create a slide focused on visual content with minimal text"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply background with theme
        self._apply_professional_background(slide, theme, is_title_slide=False)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = self._clean_markdown_formatting(slide_data.get('title', 'Visual Overview'))
        p.font.name = theme.font_scheme.title_font
        p.font.size = Pt(theme.font_scheme.title_size)
        p.font.bold = True
        p.font.color.rgb = theme.color_scheme.primary
        p.alignment = PP_ALIGN.CENTER
        
        # Large content area for image description
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(2))
        content_frame = content_box.text_frame
        
        content_text = slide_data.get('content', slide_data.get('image_description', 'Visual content will be displayed here'))
        self._add_paragraph_content(content_frame, content_text, theme)
        
        # Key points below
        if slide_data.get('bullet_points'):
            points_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(3))
            points_frame = points_box.text_frame
            self._add_bullet_points_with_icons(points_frame, slide_data['bullet_points'][:4], theme)
    
    def get_project_status(self, project_id):
        """Get the current status of a project"""
        return self.projects.get(project_id, {'status': 'not_found'})
    
    def download_project(self, project_id):
        """Prepare project for download"""
        project = self.projects.get(project_id)
        if not project or project['status'] != 'completed':
            return None
        
        return project.get('ppt_path')
    
    def list_projects(self):
        """List all projects"""
        return list(self.projects.values())
    
    def delete_project(self, project_id):
        """Delete a project and its files"""
        project = self.projects.get(project_id)
        if project:
            # Delete files
            if 'ppt_path' in project and os.path.exists(project['ppt_path']):
                os.remove(project['ppt_path'])
            
            # Remove from memory
            del self.projects[project_id]
            return True
        return False

